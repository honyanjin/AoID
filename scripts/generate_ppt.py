from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# Constants - CI Colors
COLOR_PRIMARY = RGBColor(25, 118, 210)      # #1976d2
COLOR_ACCENT = RGBColor(66, 165, 245)       # #42a5f5
COLOR_BG_GRAY = RGBColor(245, 245, 245)     # #f5f5f5
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_MAIN = RGBColor(51, 51, 51)      # #333333
COLOR_TEXT_LIGHT = RGBColor(100, 100, 100)
COLOR_GRAY = RGBColor(128, 128, 128)

FONT_MAIN = 'Malgun Gothic'
FONT_CODE = 'Consolas'

# Paths
ASSET_DIR = os.path.join("frontend", "src", "assets", "images")
# User specified correct logo
LOGO_PATH = os.path.join("AssetsFlies", "Logo", "AoID_Logo_V1.svg")

# Icon Mapping (Filename in asset_dir)
ICONS = {
    "React": "logo.png", # Fallback or specific
    "Node": "NodeJs_Logo.svg",
    "Postgres": "PostgreSQL_Logo.svg",
    "AWS": "Amazon_Web_Services_Logo.svg",
    "Docker": "Docker_Logo.svg",
    "Kubernetes": "Kubernetes_Logo.svg",
    "OpenAI": "OpenAI_Logo.svg",
    "Spring": "Spring_Framework_Logo.svg" # Just in case
}

def set_font(paragraph, size=16, bold=False, color=COLOR_TEXT_MAIN, font_name=FONT_MAIN):
    run = paragraph.add_run() if not paragraph.runs else paragraph.runs[0]
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = color
    return run

def create_card(slide, left, top, width, height, bg_color=COLOR_WHITE):
    """Creates a card background (rounded rectangle)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False # Reset
    return shape

def add_header(slide, title, logo_path=LOGO_PATH):
    # Header Background
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(8), Inches(1))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=28, bold=True, color=COLOR_WHITE)
    
    # Logo
    # Logo - Robust Handling
    logo_candidates = [
        os.path.join("frontend", "src", "assets", "images", "logo.png"), # V1 PNG (Assumed)
        os.path.join("AssetsFlies", "Logo", "AoID_Logo_V1.svg"),
        os.path.join("AssetsFlies", "Logo", "AoID_Logo_V1.ico")
        # Removed V0 PNG to avoid "Old Logo" appearing
    ]
    
    for path in logo_candidates:
        if os.path.exists(path):
            try:
                slide.shapes.add_picture(path, Inches(0.2), Inches(0.1), height=Inches(0.8))
                return
            except Exception:
                continue
    # Fallback
    print(f"Warning: Could not load logo for header: {title}")

def add_diagram_node(slide, text, x, y, w, h, color=COLOR_ACCENT, text_color=COLOR_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    set_font(p, size=14, bold=True, color=text_color)
    return shape

def connect_shapes(slide, shape1, shape2):
    # Simple direct connector (center to center implied by default args usually, but let's just place a line)
    # Getting precise connection points in python-pptx is tricky without complex math.
    # We will approximate by drawing a line from center of 1 to center of 2.
    x1 = shape1.left + shape1.width // 2
    y1 = shape1.top + shape1.height // 2
    x2 = shape2.left + shape2.width // 2
    y2 = shape2.top + shape2.height // 2
    
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = COLOR_TEXT_LIGHT
    line.line.width = Pt(2)
    # Add arrow head
    line.line.end_arrowhead_style = 'triangle' # Simple string assignment might not work, avoiding for safety unless enum
    
def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_GRAY
    
    # Hero Card
    card = create_card(slide, Inches(1), Inches(1), Inches(8), Inches(3.625))
    
    # Large Logo Center (Use add_logo logic but manually for custom size/pos? Better to reuse logic or try-except here)
    # Re-implementing simplified fallback for Title Slide specifically
    logo_candidates = [
        os.path.join("frontend", "src", "assets", "images", "logo.png"),
        os.path.join("AssetsFlies", "Logo", "AoID_Logo_V1.svg"),
        os.path.join("AssetsFlies", "Logo", "AoID_Logo_V1.ico")
    ]
    logo_placed = False
    for path in logo_candidates:
        if os.path.exists(path):
            try:
                slide.shapes.add_picture(path, Inches(4.25), Inches(1.5), width=Inches(1.5))
                logo_placed = True
                break
            except Exception:
                continue
    if not logo_placed:
        print("Warning: Could not load large logo for title.")
    
    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=40, bold=True, color=COLOR_PRIMARY).text = "AoID Platform"
    
    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_font(p2, size=20, color=COLOR_TEXT_LIGHT).text = "Association of Independent Developers"


    # --- Slide 2: Tech Stack (Grid) ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Tech Stack & Tools")
    
    # Grid config
    cols = 3
    start_x = Inches(1.5)
    start_y = Inches(1.5)
    gap_x = Inches(2.5)
    gap_y = Inches(2.0)
    
    techs = [
        ("React 19", "Frontend", ICONS["React"]),
        ("Node.js", "Backend Runtime", ICONS["Node"]),
        ("Express 5", "Web Framework", ICONS["Node"]), # Recycle Node icon
        ("PostgreSQL", "Database", ICONS["Postgres"]),
        ("AWS", "Infrastructure", ICONS["AWS"]),
        ("Docker", "Container", ICONS["Docker"])
    ]
    
    for i, (name, role, icon_file) in enumerate(techs):
        row = i // cols
        col = i % cols
        x = start_x + (col * gap_x)
        y = start_y + (row * gap_y)
        
        # Card
        create_card(slide, x, y, Inches(2.2), Inches(1.8))
        
        # Icon
        icon_path = os.path.join(ASSET_DIR, icon_file)
        if os.path.exists(icon_path):
            try:
                slide.shapes.add_picture(icon_path, x + Inches(0.85), y + Inches(0.2), width=Inches(0.5))
            except Exception as e:
                print(f"Could not load {icon_file}: {e}")
                # Fallback
                fb = slide.shapes.add_textbox(x + Inches(0.85), y + Inches(0.2), Inches(0.5), Inches(0.5))
                fb.text_frame.text = "?"
        else:
            fb = slide.shapes.add_textbox(x + Inches(0.85), y + Inches(0.2), Inches(0.5), Inches(0.5))
            fb.text_frame.text = "?"
        
        # Text
        tb = slide.shapes.add_textbox(x, y + Inches(0.8), Inches(2.2), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=16, bold=True, color=COLOR_PRIMARY).text = name
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_font(p2, size=12, color=COLOR_TEXT_LIGHT).text = role


    # --- Slide 3: Architecture (Diagram) ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, "System Architecture")
    
    # Nodes
    client = add_diagram_node(slide, "Client\n(Browser)", Inches(4), Inches(1.5), Inches(2), Inches(0.8), color=COLOR_GRAY)
    nginx = add_diagram_node(slide, "Nginx\n(Reverse Proxy)", Inches(4), Inches(2.8), Inches(2), Inches(0.8), color=COLOR_PRIMARY)
    
    fe = add_diagram_node(slide, "Frontend\n(React)", Inches(1.5), Inches(4.2), Inches(2), Inches(0.8), color=COLOR_ACCENT)
    be = add_diagram_node(slide, "Backend\n(Express)", Inches(6.5), Inches(4.2), Inches(2), Inches(0.8), color=COLOR_ACCENT)
    
    db = add_diagram_node(slide, "PostgreSQL\n(DB)", Inches(6.5), Inches(5.5), Inches(2), Inches(0.8), color=RGBColor(51, 103, 145)) # PG Color
    
    # Connections
    connect_shapes(slide, client, nginx)
    connect_shapes(slide, nginx, fe)
    connect_shapes(slide, nginx, be)
    connect_shapes(slide, be, db)


    # --- Slide 4: Key Features (Cards) ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Key Features")
    
    features = [
        ("MotionBox CMS", "Drag & Drop UI builder engine.\nJSON-based rendering."),
        ("DevCoins Economy", "Virtual currency system.\nIntegrated with Toss Payments."),
        ("Grade System", "Automatic promotion logic.\nDetailed permission control.")
    ]
    
    y = Inches(1.5)
    for title, desc in features:
        # Card
        create_card(slide, Inches(1), y, Inches(8), Inches(1.2))
        
        # Text
        tb = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.1), Inches(7.6), Inches(1))
        p = tb.text_frame.paragraphs[0]
        set_font(p, size=18, bold=True, color=COLOR_PRIMARY).text = title
        
        p2 = tb.text_frame.add_paragraph()
        set_font(p2, size=14, color=COLOR_TEXT_MAIN).text = desc
        
        y += Inches(1.4)

    # --- Slide 5: Stats ---
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Project Scale")
    
    # Big Numbers
    stats = [
        ("39", "DB Migrations"),
        ("200+", "Components"),
        ("100%", "TypeScript"),
        ("v19", "React Version")
    ]
    
    bx = Inches(0.5)
    by = Inches(2)
    gap = Inches(2.3)
    
    for num, label in stats:
        create_card(slide, bx, by, Inches(2.1), Inches(2.1), bg_color=COLOR_PRIMARY)
        
        tb = slide.shapes.add_textbox(bx, by + Inches(0.3), Inches(2.1), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=44, bold=True, color=COLOR_WHITE).text = num
        
        tb2 = slide.shapes.add_textbox(bx, by + Inches(1.2), Inches(2.1), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        set_font(p2, size=16, color=COLOR_ACCENT).text = label
        
        bx += gap

    # Save
    out = 'AoID_Presentation_v7.pptx'
    prs.save(out)
    print(f"Saved {os.path.abspath(out)}")

if __name__ == "__main__":
    create_presentation()
